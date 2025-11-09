import os
from pptx import Presentation
from PIL import Image
import pytesseract
import io
import pandas as pd
from sentence_transformers import SentenceTransformer
import cv2
import numpy as np
from typing import List, Dict, Any
import torch
from torchvision import transforms

class PPTEmbedder:
    def __init__(self, model_name='all-MiniLM-L6-v2'):
        """初始化embedding模型"""
        self.text_model = SentenceTransformer(model_name)
        # 设置OCR路径（根据你的Tesseract安装路径调整）
        # pytesseract.pytesseract.tesseract_cmd = r'C:\Program Files\Tesseract-OCR\tesseract.exe'  # Windows
        # pytesseract.pytesseract.tesseract_cmd = '/usr/local/bin/tesseract'  # Mac/Linux
        
    def extract_ppt_content(self, ppt_path: str) -> List[Dict[str, Any]]:
        """
        提取PPT中的所有内容
        返回包含文本、图片OCR文本和表格数据的列表
        """
        prs = Presentation(ppt_path)
        slides_content = []
        
        for slide_num, slide in enumerate(prs.slides):
            slide_data = {
                'slide_number': slide_num + 1,
                'text_content': [],
                'image_text': [],
                'tables': []
            }
            
            # 提取形状中的文本
            for shape in slide.shapes:
                # 文本框内容
                if hasattr(shape, "text") and shape.text.strip():
                    text = shape.text.strip()
                    if text:  # 确保不是空字符串
                        slide_data['text_content'].append(text)
                
                # 表格内容
                if shape.has_table:
                    table_data = self._extract_table(shape.table)
                    slide_data['tables'].append(table_data)
                
                # 图片内容（OCR识别）
                if shape.shape_type == 13:  # 图片类型
                    image_text = self._extract_image_text(shape)
                    if image_text:
                        slide_data['image_text'].append(image_text)
            
            slides_content.append(slide_data)
        
        return slides_content
    
    def _extract_table(self, table) -> List[List[str]]:
        """提取表格数据"""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(row_data)
        return table_data
    
    def _extract_image_text(self, shape) -> str:
        """从图片中提取文本（OCR）"""
        try:
            # 获取图片数据
            image = shape.image
            image_bytes = image.blob
            
            # 转换为PIL Image
            img = Image.open(io.BytesIO(image_bytes))
            
            # 使用OCR识别文本
            text = pytesseract.image_to_string(img, lang='chi_sim+eng')  # 支持中英文
            
            return text.strip()
        except Exception as e:
            print(f"OCR处理失败: {e}")
            return ""
    
    def preprocess_content(self, slides_content: List[Dict[str, Any]]) -> List[str]:
        """预处理内容，准备用于embedding"""
        chunks = []
        
        for slide in slides_content:
            slide_texts = []
            
            # 添加文本内容
            if slide['text_content']:
                slide_texts.extend(slide['text_content'])
            
            # 添加图片OCR文本
            if slide['image_text']:
                for img_text in slide['image_text']:
                    if img_text:
                        slide_texts.append(f"[图片内容] {img_text}")
            
            # 添加表格内容
            if slide['tables']:
                for i, table in enumerate(slide['tables']):
                    table_text = self._table_to_text(table)
                    slide_texts.append(f"[表格{i+1}] {table_text}")
            
            # 如果本页有内容，合并为一个chunk
            if slide_texts:
                chunk = f"幻灯片{slide['slide_number']}: " + " | ".join(slide_texts)
                chunks.append(chunk)
        
        return chunks
    
    def _table_to_text(self, table_data: List[List[str]]) -> str:
        """将表格数据转换为文本描述"""
        if not table_data:
            return ""
        
        # 简单的表格文本化
        table_text = "表格内容: "
        for i, row in enumerate(table_data):
            row_text = " | ".join([cell for cell in row if cell])
            if row_text:
                table_text += f"行{i+1}: {row_text}; "
        
        return table_text
    
    def create_embeddings(self, text_chunks: List[str]):
        """为文本块创建embedding"""
        embeddings = self.text_model.encode(text_chunks)
        return embeddings
    
    def process_ppt(self, ppt_path: str):
        """完整的PPT处理流程"""
        print("正在提取PPT内容...")
        slides_content = self.extract_ppt_content(ppt_path)
        
        print("正在预处理内容...")
        text_chunks = self.preprocess_content(slides_content)
        
        print("正在生成embedding...")
        embeddings = self.create_embeddings(text_chunks)
        
        return {
            'slides_content': slides_content,
            'text_chunks': text_chunks,
            'embeddings': embeddings
        }

# 3. 使用示例
if __name__ == "__main__":
    # 初始化embedder
    embedder = PPTEmbedder()
    
    # 处理PPT文件
    ppt_path = "./source/培训资料.pptx"  # 替换为你的PPT路径
    result = embedder.process_ppt(ppt_path)
    
    # 打印结果摘要
    print(f"\n处理完成!")
    print(f"幻灯片数量: {len(result['slides_content'])}")
    print(f"文本块数量: {len(result['text_chunks'])}")
    print(f"Embedding形状: {result['embeddings'].shape}")
    
    # 查看前几个文本块
    print("\n前3个文本块:")
    for i, chunk in enumerate(result['text_chunks'][:3]):
        print(f"{i+1}. {chunk[:100]}...")
    
    # 保存结果（可选）
    import numpy as np
    np.save('ppt_embeddings.npy', result['embeddings'])
    
    # 保存文本块和对应的embedding索引
    with open('ppt_chunks.txt', 'w', encoding='utf-8') as f:
        for i, chunk in enumerate(result['text_chunks']):
            f.write(f"Chunk {i}: {chunk}\n\n")